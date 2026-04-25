"""
Document generation service.

Takes LLM-produced content and formats it into:
  .docx  — Microsoft Word   (python-docx)
  .xlsx  — Microsoft Excel  (openpyxl)
  .pdf   — PDF              (reportlab)
  .pptx  — PowerPoint       (python-pptx)
  .txt   — plain text       (stdlib)
  .csv   — CSV              (stdlib)

All functions accept a `content` dict produced by the LLM and return (bytes, filename).
Sections may carry an optional "image_bytes" key (PNG bytes) that will be embedded.
"""

from __future__ import annotations

import csv
import io
import re
from datetime import datetime
from typing import Any


# ── helpers ──────────────────────────────────────────────────────────────────

def _safe_filename(title: str, ext: str) -> str:
    slug = re.sub(r"[^\w\s-]", "", title.lower()).strip()
    slug = re.sub(r"[\s_-]+", "_", slug)[:40] or "nova_document"
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    return f"{slug}_{ts}.{ext}"


# ── Word (.docx) ──────────────────────────────────────────────────────────────

def _make_docx(content: dict[str, Any]) -> bytes:
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor as DocxRGB
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    doc = Document()

    # Tighten margins
    for sec in doc.sections:
        sec.top_margin    = Inches(1.0)
        sec.bottom_margin = Inches(1.0)
        sec.left_margin   = Inches(1.25)
        sec.right_margin  = Inches(1.25)

    # ------------------------------------------------------------------
    # Use raw paragraph + explicit run formatting for ALL text so that
    # the styling is embedded at the run level and renders correctly in
    # every viewer (Word, LibreOffice, browser docx-preview, etc.) without
    # depending on built-in style inheritance.
    # ------------------------------------------------------------------

    def _add_title(text: str) -> None:
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.paragraph_format.space_before = Pt(0)
        p.paragraph_format.space_after  = Pt(4)
        run = p.add_run(text)
        run.font.name  = "Calibri"
        run.font.size  = Pt(24)
        run.font.bold  = True
        run.font.color.rgb = DocxRGB(0x1E, 0x29, 0x3B)

    def _add_rule() -> None:
        """Thin blue horizontal rule via paragraph bottom border."""
        rule_p = doc.add_paragraph()
        rule_p.paragraph_format.space_before = Pt(0)
        rule_p.paragraph_format.space_after  = Pt(16)
        pPr = rule_p._p.get_or_add_pPr()
        pBdr = OxmlElement("w:pBdr")
        bottom = OxmlElement("w:bottom")
        bottom.set(qn("w:val"), "single")
        bottom.set(qn("w:sz"), "6")
        bottom.set(qn("w:space"), "1")
        bottom.set(qn("w:color"), "2563EB")
        pBdr.append(bottom)
        pPr.append(pBdr)

    def _add_section_heading(text: str) -> None:
        p = doc.add_paragraph()
        p.paragraph_format.space_before = Pt(20)
        p.paragraph_format.space_after  = Pt(4)
        run = p.add_run(text)
        run.font.name  = "Calibri"
        run.font.size  = Pt(14)
        run.font.bold  = True
        run.font.color.rgb = DocxRGB(0x1D, 0x4E, 0xD8)

    def _add_body_para(text: str) -> None:
        p = doc.add_paragraph()
        p.paragraph_format.space_before = Pt(0)
        p.paragraph_format.space_after  = Pt(8)
        p.paragraph_format.line_spacing = Pt(18)
        run = p.add_run(text)
        run.font.name = "Calibri"
        run.font.size = Pt(11)
        run.font.color.rgb = DocxRGB(0x1E, 0x29, 0x3B)

    def _add_subtitle(text: str) -> None:
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.paragraph_format.space_after = Pt(14)
        run = p.add_run(text)
        run.font.name   = "Calibri"
        run.font.size   = Pt(12)
        run.font.italic = True
        run.font.color.rgb = DocxRGB(0x64, 0x74, 0x8B)

    # ── Title block ────────────────────────────────────────────────────
    title = (content.get("title") or "Document").strip()
    _add_title(title)
    _add_rule()

    if content.get("subtitle"):
        _add_subtitle(content["subtitle"])

    # ── Sections ───────────────────────────────────────────────────────
    for section in content.get("sections", []):
        heading = (section.get("heading") or "").strip()
        if heading:
            _add_section_heading(heading)

        for para in section.get("paragraphs", []):
            if para.strip():
                _add_body_para(para.strip())

        img_bytes = section.get("image_bytes")
        if img_bytes:
            try:
                img_stream = io.BytesIO(img_bytes)
                doc.add_picture(img_stream, width=Inches(4.5))
                last = doc.paragraphs[-1]
                last.alignment = WD_ALIGN_PARAGRAPH.CENTER
                last.paragraph_format.space_after = Pt(12)
            except Exception as e:
                print(f"[DocGen] DOCX image embed failed: {e}", flush=True)

    # ── Plain body fallback ────────────────────────────────────────────
    if not content.get("sections") and content.get("body"):
        for line in content["body"].split("\n"):
            if line.strip():
                _add_body_para(line.strip())

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ── Excel (.xlsx) ─────────────────────────────────────────────────────────────

def _make_xlsx(content: dict[str, Any]) -> bytes:  # noqa: PLR0912, PLR0914
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side, GradientFill
    from openpyxl.utils import get_column_letter
    from openpyxl.chart import BarChart, LineChart, PieChart, AreaChart, Reference

    # ── Style constants ───────────────────────────────────────────────────────
    HDR_FONT   = Font(bold=True, color="FFFFFF", size=11)
    HDR_FILL   = PatternFill("solid", fgColor="2563EB")
    HDR_ALIGN  = Alignment(horizontal="center", vertical="center", wrap_text=True)
    TITLE_FONT = Font(bold=True, size=14, color="1E293B")
    TITLE_FILL = PatternFill("solid", fgColor="EFF6FF")  # blue-50
    ALT_FILL   = PatternFill("solid", fgColor="F8FAFF")   # very light blue
    THIN_SIDE  = Side(style="thin", color="CBD5E1")
    THIN_BRD   = Border(left=THIN_SIDE, right=THIN_SIDE, top=THIN_SIDE, bottom=THIN_SIDE)

    def _coerce(val: Any) -> Any:
        """Try numeric coercion so charts work on number columns."""
        if isinstance(val, (int, float)):
            return val
        try:
            return int(val)
        except (ValueError, TypeError):
            pass
        try:
            return float(str(val).replace(",", ""))
        except (ValueError, TypeError):
            return val

    def _write_sheet(wb: Any, sheet_def: dict, title: str, first: bool) -> None:
        ws = wb.active if first else wb.create_sheet()
        ws.title = sheet_def.get("sheet_title", "Sheet1")[:31]

        headers = sheet_def.get("headers", [])
        rows    = sheet_def.get("rows", [])
        n_cols  = max(len(headers), max((len(r) for r in rows), default=1))

        row_idx = 1

        # ── Title row ─────────────────────────────────────────────────────────
        span_col = get_column_letter(max(n_cols, 1))
        ws.merge_cells(f"A{row_idx}:{span_col}{row_idx}")
        tc = ws.cell(row=row_idx, column=1, value=title)
        tc.font      = TITLE_FONT
        tc.fill      = TITLE_FILL
        tc.alignment = Alignment(horizontal="center", vertical="center")
        ws.row_dimensions[row_idx].height = 28
        row_idx += 1

        # ── Column headers ────────────────────────────────────────────────────
        header_row = row_idx
        if headers:
            for ci, hdr in enumerate(headers, start=1):
                c = ws.cell(row=row_idx, column=ci, value=hdr)
                c.font      = HDR_FONT
                c.fill      = HDR_FILL
                c.alignment = HDR_ALIGN
                c.border    = THIN_BRD
            ws.row_dimensions[row_idx].height = 22
            row_idx += 1

        # ── Data rows ─────────────────────────────────────────────────────────
        data_start = row_idx
        for ri, data_row in enumerate(rows):
            fill = ALT_FILL if ri % 2 == 1 else None
            for ci, val in enumerate(data_row, start=1):
                c = ws.cell(row=row_idx, column=ci, value=_coerce(val))
                c.border    = THIN_BRD
                c.alignment = Alignment(vertical="center")
                if fill:
                    c.fill = fill
            ws.row_dimensions[row_idx].height = 18
            row_idx += 1
        data_end = row_idx - 1

        # ── Auto-fit column widths ─────────────────────────────────────────────
        for col_cells in ws.columns:
            max_len   = 0
            col_letter = get_column_letter(col_cells[0].column)
            for cell in col_cells:
                if cell.value is not None:
                    max_len = max(max_len, len(str(cell.value)))
            ws.column_dimensions[col_letter].width = min(max_len + 4, 40)

        # ── Chart ──────────────────────────────────────────────────────────────
        chart_cfg = sheet_def.get("chart")
        if chart_cfg and rows:
            ctype      = chart_cfg.get("type", "bar").lower()
            ctitle     = chart_cfg.get("title", "")
            cat_col    = int(chart_cfg.get("category_col", 0)) + 1   # 1-indexed
            val_cols   = [int(c) + 1 for c in chart_cfg.get("value_cols", [1])]

            # Build chart object
            if ctype == "line":
                chart = LineChart()
                chart.style = 12
            elif ctype == "pie":
                chart = PieChart()
                chart.style = 10
            elif ctype == "area":
                chart = AreaChart()
                chart.style = 11
                chart.grouping = "stacked"
            else:  # default bar / column
                chart = BarChart()
                chart.type     = "col"
                chart.grouping = "clustered"
                chart.style    = 10

            chart.title   = ctitle
            chart.width   = 18
            chart.height  = 12

            if ctype == "pie":
                # Pie chart uses a single series
                vc = val_cols[0]
                data = Reference(ws, min_col=vc, max_col=vc,
                                 min_row=header_row, max_row=data_end)
                chart.add_data(data, titles_from_data=True)
                cats = Reference(ws, min_col=cat_col,
                                 min_row=data_start, max_row=data_end)
                chart.set_categories(cats)
            else:
                # Multi-series bar / line / area
                min_vc = min(val_cols)
                max_vc = max(val_cols)
                data = Reference(ws, min_col=min_vc, max_col=max_vc,
                                 min_row=header_row, max_row=data_end)
                chart.add_data(data, titles_from_data=True)
                cats = Reference(ws, min_col=cat_col,
                                 min_row=data_start, max_row=data_end)
                chart.set_categories(cats)

            # Place chart two columns after the data
            chart_anchor = f"{get_column_letter(n_cols + 2)}{header_row}"
            ws.add_chart(chart, chart_anchor)

        # ── Embedded image (placed below chart / below data table) ────────────
        img_bytes = sheet_def.get("image_bytes")
        if img_bytes:
            try:
                from openpyxl.drawing.image import Image as XLImage
                xl_img = XLImage(io.BytesIO(img_bytes))
                # Scale to reasonable size
                xl_img.width  = 360   # px
                xl_img.height = 225   # px  (16:10 ratio)
                # Anchor: below chart if chart exists, otherwise below data
                img_row = data_end + 2
                img_col = n_cols + 2 if not sheet_def.get("chart") else n_cols + 2 + 13
                ws.add_image(xl_img, f"{get_column_letter(img_col)}{img_row}")
            except Exception as e:
                print(f"[DocGen] XLSX image embed failed: {e}", flush=True)

    # ── Normalise content to sheet list ──────────────────────────────────────
    wb    = Workbook()
    title = content.get("title", "Spreadsheet")

    sheets = content.get("sheets")
    if sheets:
        # New multi-sheet format
        for i, sheet_def in enumerate(sheets):
            _write_sheet(wb, sheet_def, title, first=(i == 0))
    else:
        # Legacy flat format — single sheet
        _write_sheet(wb, content, title, first=True)

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


# ── PDF ───────────────────────────────────────────────────────────────────────

def _make_pdf(content: dict[str, Any]) -> bytes:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import mm, inch
    from reportlab.lib import colors
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, HRFlowable, Image as RLImage

    buf = io.BytesIO()
    doc = SimpleDocTemplate(
        buf,
        pagesize=A4,
        rightMargin=20 * mm,
        leftMargin=20 * mm,
        topMargin=20 * mm,
        bottomMargin=20 * mm,
    )

    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        "NovaTitle",
        parent=styles["Title"],
        fontSize=24,
        spaceAfter=4,
        spaceBefore=0,
        fontName="Helvetica-Bold",
        textColor=colors.HexColor("#1e293b"),
        alignment=1,  # CENTER
    )
    h1_style = ParagraphStyle(
        "NovaH1",
        parent=styles["Heading1"],
        fontSize=14,
        spaceBefore=16,
        spaceAfter=6,
        fontName="Helvetica-Bold",
        textColor=colors.HexColor("#2563eb"),
    )
    body_style = ParagraphStyle(
        "NovaBody",
        parent=styles["Normal"],
        fontSize=11,
        leading=18,
        spaceAfter=10,
        fontName="Helvetica",
        textColor=colors.HexColor("#1e293b"),
    )

    story = []

    title = content.get("title", "Document")
    story.append(Paragraph(title, title_style))
    story.append(HRFlowable(width="100%", thickness=1, color=colors.HexColor("#2563eb")))
    story.append(Spacer(1, 8 * mm))

    if content.get("subtitle"):
        story.append(Paragraph(content["subtitle"], styles["Italic"]))
        story.append(Spacer(1, 4 * mm))

    for section in content.get("sections", []):
        if section.get("heading"):
            story.append(Paragraph(section["heading"], h1_style))

        # Body paragraphs first — image comes after the text it illustrates
        for para in section.get("paragraphs", []):
            safe = para.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            story.append(Paragraph(safe, body_style))

        # Embedded image after the section's text
        img_bytes = section.get("image_bytes")
        if img_bytes:
            try:
                img_stream = io.BytesIO(img_bytes)
                img = RLImage(img_stream, width=4.5 * inch, height=3.0 * inch)
                img.hAlign = "CENTER"
                story.append(Spacer(1, 4 * mm))
                story.append(img)
                story.append(Spacer(1, 6 * mm))
            except Exception as e:
                print(f"[DocGen] PDF image embed failed: {e}", flush=True)

    if not content.get("sections") and content.get("body"):
        for line in content["body"].split("\n"):
            if line.strip():
                safe = line.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
                story.append(Paragraph(safe, body_style))

    doc.build(story)
    return buf.getvalue()


# ── PowerPoint (.pptx) ────────────────────────────────────────────────────────

def _make_pptx(content: dict[str, Any]) -> bytes:  # noqa: PLR0914
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    # ── Palette ──────────────────────────────────────────────────────────────
    ACCENT    = RGBColor(0x26, 0x63, 0xEB)   # blue-600
    ACCENT2   = RGBColor(0x7C, 0x3A, 0xED)   # violet-600
    DARK_BG   = RGBColor(0x0F, 0x17, 0x2A)   # slate-900
    LIGHT_BG  = RGBColor(0xF1, 0xF5, 0xFD)   # blue-50
    WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
    DARK_TXT  = RGBColor(0x1E, 0x29, 0x3B)   # slate-800
    GRAY_TXT  = RGBColor(0x94, 0xA3, 0xB8)   # slate-400

    W = Inches(13.33)
    H = Inches(7.5)

    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]  # completely blank

    # ── Helper: solid background fill ────────────────────────────────────────
    def _fill_bg(slide, rgb: RGBColor) -> None:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = rgb

    # ── Helper: add a filled rectangle ───────────────────────────────────────
    def _rect(slide, left, top, width, height, rgb: RGBColor, line: bool = False):
        shp = slide.shapes.add_shape(1, left, top, width, height)
        shp.fill.solid()
        shp.fill.fore_color.rgb = rgb
        if line:
            shp.line.color.rgb = rgb
        else:
            shp.line.fill.background()
        return shp

    # ── Helper: add a text box ────────────────────────────────────────────────
    def _label(
        slide, text: str, left, top, width, height,
        size: int, bold: bool = False, italic: bool = False,
        color: RGBColor = WHITE, align=PP_ALIGN.LEFT, wrap: bool = True,
    ):
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = wrap
        p  = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text          = text
        run.font.size     = Pt(size)
        run.font.bold     = bold
        run.font.italic   = italic
        run.font.color.rgb = color
        return tb

    sections    = content.get("sections", [])
    total_slides = 1 + len(sections)

    # ── Slide 1: Title ────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank)
    _fill_bg(slide, DARK_BG)

    # Thin accent bar on left edge
    _rect(slide, 0, 0, Inches(0.08), H, ACCENT)

    # Optional title-slide background image (first section's image if any)
    title_img = content.get("title_image_bytes") or (
        sections[0].get("image_bytes") if sections else None
    )
    if title_img:
        try:
            img_stream = io.BytesIO(title_img)
            pic = slide.shapes.add_picture(img_stream, Inches(7.5), Inches(0.5), Inches(5.5), Inches(6.5))
            # Semi-transparent dark overlay rectangle so text stays readable
            ovl = _rect(slide, Inches(7.5), Inches(0.5), Inches(5.5), Inches(6.5),
                        RGBColor(0x0F, 0x17, 0x2A))
            # Set to 50 % transparency via XML
            try:
                from pptx.oxml.ns import qn
                from lxml import etree
                spPr = ovl._element.spPr
                solidFill = spPr.find(qn("a:solidFill"))
                if solidFill is not None:
                    srgb = solidFill.find(qn("a:srgbClr"))
                    if srgb is None:
                        srgb = etree.SubElement(solidFill, qn("a:srgbClr"))
                        srgb.set("val", "0F172A")
                    alpha = etree.SubElement(srgb, qn("a:alpha"))
                    alpha.set("val", "70000")  # 70% opaque → 30% transparent
            except Exception:
                pass
        except Exception as e:
            print(f"[DocGen] PPTX title bg image failed: {e}", flush=True)

    # Main title text
    _label(
        slide, content.get("title", "Presentation"),
        Inches(0.5), Inches(2.0), Inches(12.33), Inches(2.5),
        size=44, bold=True, align=PP_ALIGN.CENTER,
    )

    # Decorative divider line
    _rect(slide, Inches(2.0), Inches(4.55), Inches(9.33), Pt(2), ACCENT)

    if content.get("subtitle"):
        _label(
            slide, content["subtitle"],
            Inches(0.5), Inches(4.75), Inches(12.33), Inches(1.0),
            size=20, italic=True, color=GRAY_TXT, align=PP_ALIGN.CENTER,
        )

    _label(slide, f"1 / {total_slides}",
           Inches(12.3), Inches(7.15), Inches(1.0), Inches(0.3),
           size=9, color=GRAY_TXT, align=PP_ALIGN.RIGHT)

    # ── Content slides ────────────────────────────────────────────────────────
    HEADER_H = Inches(1.15)
    PAD      = Inches(0.2)
    BODY_TOP = HEADER_H + PAD
    BODY_H   = H - BODY_TOP - Inches(0.4)

    for idx, sec in enumerate(sections, start=2):
        slide = prs.slides.add_slide(blank)
        _fill_bg(slide, LIGHT_BG)

        # Gradient-look header: two overlapping rects
        _rect(slide, 0, 0, W, HEADER_H, ACCENT)
        _rect(slide, W - Inches(3.0), 0, Inches(3.0), HEADER_H, ACCENT2)

        # Slide heading
        _label(
            slide, sec.get("heading", ""),
            Inches(0.3), Inches(0.12), Inches(12.7), HEADER_H - Inches(0.12),
            size=24, bold=True, color=WHITE,
        )

        # ── Image placement ───────────────────────────────────────────────────
        img_bytes = sec.get("image_bytes")
        text_w    = W - Inches(0.6)   # default: full text width

        if img_bytes:
            IMG_W = Inches(5.5)
            IMG_X = W - IMG_W - Inches(0.15)
            text_w = IMG_X - Inches(0.15)
            try:
                img_stream = io.BytesIO(img_bytes)
                # White card behind image
                _rect(slide, IMG_X - Pt(6), BODY_TOP - Pt(4),
                      IMG_W + Pt(12), BODY_H + Pt(8), WHITE)
                slide.shapes.add_picture(
                    img_stream, IMG_X, BODY_TOP + PAD / 2,
                    IMG_W, BODY_H - PAD,
                )
            except Exception as e:
                print(f"[DocGen] PPTX image embed failed: {e}", flush=True)
                text_w = W - Inches(0.6)

        # ── Bullet text ───────────────────────────────────────────────────────
        tb = slide.shapes.add_textbox(Inches(0.3), BODY_TOP + PAD / 2, text_w, BODY_H - PAD)
        tf = tb.text_frame
        tf.word_wrap = True

        for j, para in enumerate(sec.get("paragraphs", [])):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.level        = 0
            p.space_before = Pt(8)
            # Bullet dot
            r0 = p.add_run()
            r0.text            = "● "
            r0.font.color.rgb  = ACCENT
            r0.font.size       = Pt(15)
            r0.font.bold       = True
            # Body text
            r1 = p.add_run()
            r1.text            = para
            r1.font.color.rgb  = DARK_TXT
            r1.font.size       = Pt(15)

        # Slide number
        _label(slide, f"{idx} / {total_slides}",
               Inches(12.3), Inches(7.15), Inches(1.0), Inches(0.3),
               size=9, color=GRAY_TXT, align=PP_ALIGN.RIGHT)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── Plain text ────────────────────────────────────────────────────────────────

def _make_txt(content: dict[str, Any]) -> bytes:
    lines: list[str] = []
    title = content.get("title", "")
    if title:
        lines.append(title.upper())
        lines.append("=" * len(title))
        lines.append("")

    for section in content.get("sections", []):
        if section.get("heading"):
            heading = section["heading"]
            lines.append(heading)
            lines.append("-" * len(heading))
        for para in section.get("paragraphs", []):
            lines.append(para)
            lines.append("")

    if not content.get("sections") and content.get("body"):
        lines.append(content["body"])

    return "\n".join(lines).encode("utf-8")


# ── CSV ───────────────────────────────────────────────────────────────────────

def _make_csv(content: dict[str, Any]) -> bytes:
    buf = io.StringIO()
    writer = csv.writer(buf)

    headers = content.get("headers", [])
    if headers:
        writer.writerow(headers)

    for row in content.get("rows", []):
        writer.writerow(row)

    return buf.getvalue().encode("utf-8")


# ── Public API ────────────────────────────────────────────────────────────────

_MAKERS = {
    "docx": (_make_docx, "application/vnd.openxmlformats-officedocument.wordprocessingml.document"),
    "xlsx": (_make_xlsx, "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"),
    "pdf":  (_make_pdf,  "application/pdf"),
    "pptx": (_make_pptx, "application/vnd.openxmlformats-officedocument.presentationml.presentation"),
    "txt":  (_make_txt,  "text/plain; charset=utf-8"),
    "csv":  (_make_csv,  "text/csv; charset=utf-8"),
}


def generate_document(
    content: dict[str, Any],
    fmt: str,
) -> tuple[bytes, str, str]:
    """
    Convert *content* dict to a formatted document.

    Returns (file_bytes, filename, content_type).

    content dict schema (all fields optional except one of sections/body/rows):
    {
        "title":     str,
        "subtitle":  str,
        "sections": [{"heading": str, "paragraphs": [str, ...]}, ...],
        "body":      str,          # plain text fallback for docx/pdf/txt
        "headers":  [str, ...],    # for xlsx/csv
        "rows":    [[...], ...],   # for xlsx/csv
        "sheet_title": str,        # for xlsx
    }
    """
    fmt = fmt.lower().lstrip(".")
    if fmt not in _MAKERS:
        raise ValueError(f"Unsupported format: {fmt!r}. Choose from: {', '.join(_MAKERS)}")

    maker_fn, mime = _MAKERS[fmt]
    title = content.get("title", "nova_document")
    filename = _safe_filename(title, fmt)

    print(f"[DocGen] Building {filename}…", flush=True)
    file_bytes = maker_fn(content)
    print(f"[DocGen] Done — {len(file_bytes) // 1024} KB {fmt.upper()}", flush=True)

    return file_bytes, filename, mime
